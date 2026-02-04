"""PPTX 내보내기 모듈"""

from pathlib import Path
from typing import Optional

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

from src.core.presentation import Presentation, Slide, SlideLayoutType
from src.core.themes import Theme, get_theme


# 슬라이드 크기 (16:9)
SLIDE_WIDTH = Inches(13.333)  # 1280px at 96dpi
SLIDE_HEIGHT = Inches(7.5)    # 720px at 96dpi


def _hex_to_rgb(hex_color: str) -> str:
    """HEX 색상을 RGB 문자열로 변환 (# 제거)"""
    return hex_color.lstrip("#")


class PptxExporter:
    """PPTX 파일 내보내기"""

    def __init__(self, theme: Optional[Theme] = None):
        self.pptx = PptxPresentation()
        self.pptx.slide_width = SLIDE_WIDTH
        self.pptx.slide_height = SLIDE_HEIGHT
        self.theme = theme or get_theme("default")

    def export(self, presentation: Presentation, output_path: Path) -> bool:
        """프레젠테이션을 PPTX 파일로 내보내기"""
        try:
            # 프레젠테이션의 테마 설정 사용
            if presentation.theme:
                self.theme = get_theme(presentation.theme)

            for slide in presentation.slides:
                self._add_slide(slide)

            self.pptx.save(str(output_path))
            return True
        except Exception as e:
            print(f"PPTX 내보내기 실패: {e}")
            return False

    def _add_slide(self, slide: Slide):
        """슬라이드 추가"""
        # 빈 레이아웃 사용 (인덱스 6)
        blank_layout = self.pptx.slide_layouts[6]
        pptx_slide = self.pptx.slides.add_slide(blank_layout)

        # 배경색 설정
        if slide.background_color:
            self._set_background_color(pptx_slide, slide.background_color)

        # 레이아웃에 따른 콘텐츠 추가
        match slide.layout:
            case SlideLayoutType.TITLE:
                self._add_title_slide(pptx_slide, slide)
            case SlideLayoutType.TITLE_CONTENT:
                self._add_title_content_slide(pptx_slide, slide)
            case SlideLayoutType.BULLET_POINTS:
                self._add_bullet_slide(pptx_slide, slide)
            case SlideLayoutType.TITLE_IMAGE:
                self._add_title_image_slide(pptx_slide, slide)
            case _:
                self._add_title_content_slide(pptx_slide, slide)

        # 발표자 노트 추가
        if slide.notes:
            pptx_slide.notes_slide.notes_text_frame.text = slide.notes

    def _set_background_color(self, pptx_slide, color: str):
        """슬라이드 배경색 설정"""
        background = pptx_slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor.from_string(color)

    def _add_title_slide(self, pptx_slide, slide: Slide):
        """제목 슬라이드 추가"""
        colors = self.theme.colors

        # 제목
        title_box = pptx_slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor.from_string(_hex_to_rgb(colors.text_primary))
        title_para.alignment = PP_ALIGN.CENTER

        # 악센트 바 (테마 색상)
        accent_bar = pptx_slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(4), Inches(4.0), Inches(5.33), Inches(0.05)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = RGBColor.from_string(_hex_to_rgb(colors.primary))
        accent_bar.line.fill.background()

        # 부제목
        if slide.subtitle:
            subtitle_box = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(4.2), Inches(12.33), Inches(0.8)
            )
            subtitle_frame = subtitle_box.text_frame
            subtitle_para = subtitle_frame.paragraphs[0]
            subtitle_para.text = slide.subtitle
            subtitle_para.font.size = Pt(24)
            subtitle_para.font.color.rgb = RGBColor.from_string(_hex_to_rgb(colors.text_secondary))
            subtitle_para.alignment = PP_ALIGN.CENTER

    def _add_title_content_slide(self, pptx_slide, slide: Slide):
        """제목 + 내용 슬라이드 추가"""
        colors = self.theme.colors

        # 제목
        title_box = pptx_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor.from_string(_hex_to_rgb(colors.text_primary))

        # 제목 아래 악센트 라인
        line = pptx_slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(_hex_to_rgb(colors.primary))
        line.line.fill.background()

        # 내용
        if slide.content:
            content_box = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            content_para = content_frame.paragraphs[0]
            content_para.text = slide.content
            content_para.font.size = Pt(18)
            content_para.font.color.rgb = RGBColor.from_string(_hex_to_rgb(colors.text_secondary))

    def _add_bullet_slide(self, pptx_slide, slide: Slide):
        """글머리 기호 슬라이드 추가"""
        colors = self.theme.colors

        # 제목
        title_box = pptx_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor.from_string(_hex_to_rgb(colors.text_primary))

        # 제목 아래 악센트 라인
        line = pptx_slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(_hex_to_rgb(colors.primary))
        line.line.fill.background()

        # 글머리 기호 목록
        if slide.bullet_points:
            content_box = pptx_slide.shapes.add_textbox(
                Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5)
            )
            content_frame = content_box.text_frame
            content_frame.word_wrap = True

            for i, bullet in enumerate(slide.bullet_points):
                if i == 0:
                    para = content_frame.paragraphs[0]
                else:
                    para = content_frame.add_paragraph()

                para.text = f"• {bullet}"
                para.font.size = Pt(18)
                para.font.color.rgb = RGBColor.from_string(_hex_to_rgb(colors.text_secondary))
                para.space_after = Pt(12)

    def _add_title_image_slide(self, pptx_slide, slide: Slide):
        """제목 + 이미지 슬라이드 추가"""
        colors = self.theme.colors

        # 제목
        title_box = pptx_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.4), Inches(12.33), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_para = title_frame.paragraphs[0]
        title_para.text = slide.title
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor.from_string(_hex_to_rgb(colors.text_primary))

        # 제목 아래 악센트 라인
        line = pptx_slide.shapes.add_shape(
            1,  # MSO_SHAPE.RECTANGLE
            Inches(0.5), Inches(1.1), Inches(12.33), Inches(0.03)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor.from_string(_hex_to_rgb(colors.primary))
        line.line.fill.background()

        # 이미지
        if slide.image_url and Path(slide.image_url).exists():
            pptx_slide.shapes.add_picture(
                slide.image_url,
                Inches(2), Inches(1.5),
                width=Inches(9.33), height=Inches(5.5)
            )


def export_to_pptx(presentation: Presentation, output_path: Path) -> bool:
    """프레젠테이션을 PPTX로 내보내기 (편의 함수)"""
    exporter = PptxExporter()
    return exporter.export(presentation, output_path)

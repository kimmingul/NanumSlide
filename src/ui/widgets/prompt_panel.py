"""AI 프롬프트 패널 위젯"""

from PySide6.QtWidgets import (
    QWidget,
    QVBoxLayout,
    QHBoxLayout,
    QLabel,
    QTextEdit,
    QPushButton,
    QComboBox,
    QProgressBar,
    QFrame,
    QScrollArea,
    QGroupBox,
    QSpinBox,
    QFileDialog,
    QMessageBox,
)
from PySide6.QtCore import Qt, Signal
from pathlib import Path

from src.core.themes import get_theme_names, get_theme_by_display_name, Theme
from src.services.llm_client import get_available_models


class PromptPanel(QWidget):
    """AI 프롬프트 입력 및 생성 패널"""

    generation_requested = Signal(str)  # 생성 요청 시그널 (프롬프트)
    generation_cancelled = Signal()  # 생성 취소 시그널
    theme_changed = Signal(object)  # 테마 변경 시그널 (Theme 객체)

    def __init__(self):
        super().__init__()
        self.is_generating = False
        self.uploaded_files: list[Path] = []
        self._setup_ui()

    def _setup_ui(self):
        """UI 구성"""
        layout = QVBoxLayout(self)
        layout.setContentsMargins(12, 12, 12, 12)
        layout.setSpacing(12)

        # 헤더
        header = QLabel("AI 프레젠테이션 생성")
        header.setStyleSheet("font-weight: bold; font-size: 14px;")
        layout.addWidget(header)

        # 스크롤 영역
        scroll_area = QScrollArea()
        scroll_area.setWidgetResizable(True)
        scroll_area.setFrameShape(QFrame.Shape.NoFrame)

        scroll_content = QWidget()
        scroll_layout = QVBoxLayout(scroll_content)
        scroll_layout.setContentsMargins(0, 0, 0, 0)
        scroll_layout.setSpacing(12)

        # === 프롬프트 입력 ===
        prompt_group = QGroupBox("프롬프트")
        prompt_layout = QVBoxLayout(prompt_group)

        self.prompt_input = QTextEdit()
        self.prompt_input.setPlaceholderText(
            "프레젠테이션 주제를 입력하세요.\n\n"
            "예시:\n"
            "- 2024년 마케팅 전략 발표자료\n"
            "- Python 기초 강의 자료\n"
            "- 스타트업 투자 피칭 덱"
        )
        self.prompt_input.setMinimumHeight(150)
        prompt_layout.addWidget(self.prompt_input)

        scroll_layout.addWidget(prompt_group)

        # === 옵션 설정 ===
        options_group = QGroupBox("생성 옵션")
        options_layout = QVBoxLayout(options_group)

        # 슬라이드 수
        slide_count_layout = QHBoxLayout()
        slide_count_layout.addWidget(QLabel("슬라이드 수:"))
        self.slide_count_spin = QSpinBox()
        self.slide_count_spin.setRange(3, 30)
        self.slide_count_spin.setValue(8)
        slide_count_layout.addWidget(self.slide_count_spin)
        slide_count_layout.addStretch()
        options_layout.addLayout(slide_count_layout)

        # AI 모델 선택
        model_layout = QHBoxLayout()
        model_layout.addWidget(QLabel("AI 모델:"))
        self.model_combo = QComboBox()
        self.model_combo.setMinimumWidth(280)
        self.refresh_models()  # 초기 모델 목록 로드
        model_layout.addWidget(self.model_combo)

        # 새로고침 버튼
        refresh_btn = QPushButton("↻")
        refresh_btn.setFixedSize(28, 28)
        refresh_btn.setToolTip("모델 목록 새로고침")
        refresh_btn.clicked.connect(self.refresh_models)
        model_layout.addWidget(refresh_btn)

        options_layout.addLayout(model_layout)

        # 언어 선택
        lang_layout = QHBoxLayout()
        lang_layout.addWidget(QLabel("언어:"))
        self.lang_combo = QComboBox()
        self.lang_combo.addItems(["한국어", "English", "日本語", "中文"])
        lang_layout.addWidget(self.lang_combo)
        options_layout.addLayout(lang_layout)

        scroll_layout.addWidget(options_group)

        # === 템플릿 설정 (새 시스템) ===
        template_group = QGroupBox("템플릿")
        template_group_layout = QVBoxLayout(template_group)

        # 템플릿 선택
        template_select_layout = QHBoxLayout()
        template_select_layout.addWidget(QLabel("템플릿:"))
        self.new_template_combo = QComboBox()
        self.new_template_combo.addItems([
            "자동 선택",
            "Pitch Deck (투자 유치)",
            "Quarterly Report (보고서)",
            "Lecture (강의)",
            "Product Launch (출시)",
            "Clean Minimal (미니멀)"
        ])
        self.new_template_combo.currentTextChanged.connect(self._on_new_template_changed)
        template_select_layout.addWidget(self.new_template_combo)
        template_group_layout.addLayout(template_select_layout)

        # 템플릿 미리보기
        self.template_preview_frame = QFrame()
        self.template_preview_frame.setFixedHeight(60)
        self.template_preview_frame.setStyleSheet(
            "background: qlineargradient(x1:0, y1:0, x2:1, y2:1, "
            "stop:0 #4a5568, stop:1 #718096); border-radius: 6px;"
        )
        template_group_layout.addWidget(self.template_preview_frame)

        # 템플릿 정보
        self.template_info_label = QLabel("AI가 주제에 맞는 템플릿을 선택합니다")
        self.template_info_label.setStyleSheet("font-size: 10px; color: gray;")
        self.template_info_label.setWordWrap(True)
        template_group_layout.addWidget(self.template_info_label)

        scroll_layout.addWidget(template_group)

        # === 테마/색상 설정 ===
        theme_group = QGroupBox("색상 테마")
        theme_layout = QVBoxLayout(theme_group)

        # 테마 선택
        color_layout = QHBoxLayout()
        color_layout.addWidget(QLabel("스타일:"))
        self.template_combo = QComboBox()
        self.template_combo.addItems(get_theme_names())
        self.template_combo.currentTextChanged.connect(self._on_theme_changed)
        color_layout.addWidget(self.template_combo)
        theme_layout.addLayout(color_layout)

        # 테마 미리보기
        self.theme_preview = QFrame()
        self.theme_preview.setFixedHeight(40)
        self.theme_preview.setStyleSheet(
            "background-color: #007acc; border-radius: 4px;"
        )
        theme_layout.addWidget(self.theme_preview)

        # 색상 정보
        self.theme_colors_label = QLabel("주 색상: #007acc")
        self.theme_colors_label.setStyleSheet("font-size: 11px; opacity: 0.7;")
        theme_layout.addWidget(self.theme_colors_label)

        scroll_layout.addWidget(theme_group)

        # === 파일 업로드 ===
        upload_group = QGroupBox("참고 자료 (선택)")
        upload_layout = QVBoxLayout(upload_group)

        upload_btn_layout = QHBoxLayout()
        upload_btn = QPushButton("파일 업로드")
        upload_btn.setToolTip("PDF, DOCX, TXT, PPTX 파일을 업로드하여 참고 자료로 사용")
        upload_btn.clicked.connect(self._on_upload_clicked)
        upload_btn_layout.addWidget(upload_btn)

        clear_btn = QPushButton("초기화")
        clear_btn.setToolTip("업로드된 파일 모두 삭제")
        clear_btn.clicked.connect(self._clear_uploaded_files)
        upload_btn_layout.addWidget(clear_btn)

        upload_layout.addLayout(upload_btn_layout)

        self.uploaded_files_label = QLabel("업로드된 파일 없음")
        self.uploaded_files_label.setStyleSheet("font-size: 11px; color: gray;")
        self.uploaded_files_label.setWordWrap(True)
        upload_layout.addWidget(self.uploaded_files_label)

        scroll_layout.addWidget(upload_group)

        scroll_layout.addStretch()

        scroll_area.setWidget(scroll_content)
        layout.addWidget(scroll_area)

        # === 진행률 표시 ===
        self.progress_frame = QFrame()
        progress_layout = QVBoxLayout(self.progress_frame)
        progress_layout.setContentsMargins(0, 0, 0, 0)

        self.progress_label = QLabel("생성 중...")
        progress_layout.addWidget(self.progress_label)

        self.progress_bar = QProgressBar()
        self.progress_bar.setRange(0, 0)  # Indeterminate
        progress_layout.addWidget(self.progress_bar)

        self.progress_frame.hide()
        layout.addWidget(self.progress_frame)

        # === 생성 버튼 ===
        button_layout = QHBoxLayout()

        self.generate_btn = QPushButton("프레젠테이션 생성")
        self.generate_btn.setMinimumHeight(48)
        self.generate_btn.setStyleSheet(
            """
            QPushButton {
                background-color: #007acc;
                color: white;
                border: none;
                border-radius: 6px;
                font-weight: bold;
                font-size: 14px;
            }
            QPushButton:hover {
                background-color: #005a9e;
            }
            QPushButton:pressed {
                background-color: #004578;
            }
            QPushButton:disabled {
                background-color: #ccc;
            }
        """
        )
        self.generate_btn.clicked.connect(self._on_generate_clicked)
        button_layout.addWidget(self.generate_btn)

        layout.addLayout(button_layout)

    def _on_new_template_changed(self, template_name: str):
        """새 템플릿 시스템 - 템플릿 변경 처리"""
        templates_info = {
            "자동 선택": {
                "gradient": "#4a5568, #718096",
                "info": "AI가 주제에 맞는 템플릿을 선택합니다"
            },
            "Pitch Deck (투자 유치)": {
                "gradient": "#1a365d, #2c5282",
                "info": "투자 유치용 | 12장 권장"
            },
            "Quarterly Report (보고서)": {
                "gradient": "#0f4c81, #1e6eb8",
                "info": "비즈니스 보고서 | 15장 권장"
            },
            "Lecture (강의)": {
                "gradient": "#2d6a4f, #40916c",
                "info": "교육/강의용 | 20장 권장"
            },
            "Product Launch (출시)": {
                "gradient": "#e63946, #f72585",
                "info": "제품 출시/마케팅 | 12장 권장"
            },
            "Clean Minimal (미니멀)": {
                "gradient": "#374151, #4b5563",
                "info": "미니멀 디자인 | 10장 권장"
            }
        }

        info = templates_info.get(template_name, templates_info["자동 선택"])
        colors = info['gradient'].split(', ')
        self.template_preview_frame.setStyleSheet(
            f"background: qlineargradient(x1:0, y1:0, x2:1, y2:1, "
            f"stop:0 {colors[0]}, stop:1 {colors[1]}); border-radius: 6px;"
        )
        self.template_info_label.setText(info["info"])

    def _on_theme_changed(self, theme_name: str):
        """테마 변경 처리"""
        theme = get_theme_by_display_name(theme_name)

        # 미리보기 업데이트
        self.theme_preview.setStyleSheet(
            f"background-color: {theme.colors.primary}; border-radius: 4px;"
        )
        self.theme_colors_label.setText(f"주 색상: {theme.colors.primary}")

        # 생성 버튼 색상도 업데이트
        self.generate_btn.setStyleSheet(
            f"""
            QPushButton {{
                background-color: {theme.colors.primary};
                color: white;
                border: none;
                border-radius: 4px;
                font-weight: bold;
                font-size: 13px;
            }}
            QPushButton:hover {{
                background-color: {theme.colors.secondary};
            }}
            QPushButton:pressed {{
                background-color: {theme.colors.accent};
            }}
            QPushButton:disabled {{
                background-color: #ccc;
            }}
        """
        )

        # 테마 변경 시그널 발생
        self.theme_changed.emit(theme)

    def _on_generate_clicked(self):
        """생성 버튼 클릭 처리"""
        if self.is_generating:
            self._cancel_generation()
        else:
            self._start_generation()

    def _start_generation(self):
        """생성 시작"""
        prompt = self.prompt_input.toPlainText().strip()
        if not prompt:
            return

        self.is_generating = True
        self.generate_btn.setText("취소")
        self.progress_frame.show()
        self.progress_label.setText("AI가 프레젠테이션을 생성하고 있습니다...")

        self.generation_requested.emit(prompt)

    def _cancel_generation(self):
        """생성 취소"""
        self.is_generating = False
        self.generate_btn.setText("프레젠테이션 생성")
        self.progress_frame.hide()

        self.generation_cancelled.emit()

    def set_progress(self, message: str, value: int = -1):
        """진행률 업데이트"""
        self.progress_label.setText(message)
        if value >= 0:
            self.progress_bar.setRange(0, 100)
            self.progress_bar.setValue(value)
        else:
            self.progress_bar.setRange(0, 0)

    def generation_complete(self):
        """생성 완료"""
        self.is_generating = False
        self.generate_btn.setText("프레젠테이션 생성")
        self.progress_frame.hide()

    def focus_prompt_input(self):
        """프롬프트 입력에 포커스"""
        self.prompt_input.setFocus()

    def get_options(self) -> dict:
        """현재 설정된 옵션 반환"""
        # 템플릿 ID 매핑
        template_map = {
            "자동 선택": "auto",
            "Pitch Deck (투자 유치)": "pitch_deck",
            "Quarterly Report (보고서)": "quarterly_report",
            "Lecture (강의)": "lecture",
            "Product Launch (출시)": "product_launch",
            "Clean Minimal (미니멀)": "clean",
        }
        selected_template = self.new_template_combo.currentText()

        return {
            "slide_count": self.slide_count_spin.value(),
            "model": self.model_combo.currentText(),
            "language": self.lang_combo.currentText(),
            "template": self.template_combo.currentText(),  # 색상 테마
            "template_id": template_map.get(selected_template, "auto"),  # 새 템플릿 ID
            "reference_content": self.get_uploaded_files_content(),
        }

    def get_current_theme(self) -> Theme:
        """현재 선택된 테마 반환"""
        return get_theme_by_display_name(self.template_combo.currentText())

    def refresh_models(self):
        """모델 목록 새로고침"""
        current = self.model_combo.currentText()
        self.model_combo.clear()
        models = get_available_models()
        self.model_combo.addItems(models)

        # 기존 선택 복원
        if current and current in models:
            self.model_combo.setCurrentText(current)
        elif models and not models[0].startswith("("):
            self.model_combo.setCurrentIndex(0)

    def _on_upload_clicked(self):
        """파일 업로드 버튼 클릭"""
        file_filter = "문서 파일 (*.pdf *.docx *.doc *.txt *.pptx *.ppt *.md);;모든 파일 (*.*)"
        files, _ = QFileDialog.getOpenFileNames(
            self,
            "참고 자료 업로드",
            "",
            file_filter,
        )

        if files:
            for file_path in files:
                path = Path(file_path)
                if path not in self.uploaded_files:
                    self.uploaded_files.append(path)

            self._update_uploaded_files_label()

    def _clear_uploaded_files(self):
        """업로드된 파일 초기화"""
        self.uploaded_files.clear()
        self._update_uploaded_files_label()

    def _update_uploaded_files_label(self):
        """업로드된 파일 레이블 업데이트"""
        if not self.uploaded_files:
            self.uploaded_files_label.setText("업로드된 파일 없음")
            self.uploaded_files_label.setStyleSheet("font-size: 11px; color: gray;")
        else:
            file_names = [f.name for f in self.uploaded_files]
            self.uploaded_files_label.setText(f"파일 {len(file_names)}개: {', '.join(file_names)}")
            self.uploaded_files_label.setStyleSheet("font-size: 11px; color: #007acc;")

    def get_uploaded_files(self) -> list[Path]:
        """업로드된 파일 목록 반환"""
        return self.uploaded_files.copy()

    def get_uploaded_files_content(self) -> str:
        """업로드된 파일 내용을 텍스트로 반환"""
        if not self.uploaded_files:
            return ""

        contents = []
        for file_path in self.uploaded_files:
            try:
                content = self._read_file_content(file_path)
                if content:
                    contents.append(f"=== {file_path.name} ===\n{content}")
            except Exception as e:
                print(f"파일 읽기 실패 ({file_path.name}): {e}")

        return "\n\n".join(contents)

    def _read_file_content(self, file_path: Path) -> str:
        """파일 내용 읽기"""
        suffix = file_path.suffix.lower()

        if suffix == ".txt" or suffix == ".md":
            return file_path.read_text(encoding="utf-8")

        elif suffix == ".pdf":
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(file_path)
                text = ""
                for page in doc:
                    text += page.get_text()
                doc.close()
                return text
            except ImportError:
                QMessageBox.warning(
                    self,
                    "모듈 필요",
                    "PDF 파일을 읽으려면 PyMuPDF가 필요합니다.\npip install pymupdf",
                )
                return ""

        elif suffix in [".docx", ".doc"]:
            try:
                from docx import Document
                doc = Document(file_path)
                text = "\n".join([para.text for para in doc.paragraphs])
                return text
            except ImportError:
                QMessageBox.warning(
                    self,
                    "모듈 필요",
                    "Word 파일을 읽으려면 python-docx가 필요합니다.\npip install python-docx",
                )
                return ""

        elif suffix in [".pptx", ".ppt"]:
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                texts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            texts.append(shape.text)
                return "\n".join(texts)
            except ImportError:
                QMessageBox.warning(
                    self,
                    "모듈 필요",
                    "PowerPoint 파일을 읽으려면 python-pptx가 필요합니다.\npip install python-pptx",
                )
                return ""

        return ""
